
import os
import shutil
import tempfile
import streamlit as st
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches


# 将 PDF 转换为图像
def pdf_to_images(pdf_path):
    """将 PDF 转换为图像"""
    images = convert_from_path(pdf_path, dpi=300)  # 300 DPI 提供较高的图像质量
    return images


# 将图像插入到 PPT 中
def images_to_ppt(images, ppt_path):
    """将图像插入到 PowerPoint 文件"""
    prs = Presentation()

    for image in images:
        # 创建临时文件来保存图像
        temp_image_path = tempfile.mktemp(suffix='.png')
        image.save(temp_image_path, 'PNG')  # 保存图像为 PNG 格式

        # 添加空白幻灯片并将图像插入其中
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 使用空白布局
        slide.shapes.add_picture(temp_image_path, Inches(0), Inches(0), width=Inches(10),
                                 height=Inches(7.5))  # 将图像插入到幻灯片

        # 删除临时图像文件
        os.remove(temp_image_path)

    prs.save(ppt_path)


# 删除临时文件夹及其中内容的函数
def delete_temp_files(temp_dir):
    if os.path.exists(temp_dir):
        shutil.rmtree(temp_dir)


# Streamlit 界面
def main():
    st.title("PDF 转换为 PPT 工具")
    st.write("上传一个 PDF 文件，将其每一页转换为图像，并插入到 PowerPoint 中。")

    # 上传 PDF 文件
    uploaded_file = st.file_uploader("选择一个 PDF 文件", type=["pdf"])

    # 创建临时目录存放文件和存储转换结果
    temp_dir = None
    pdf_path = None
    ppt_path = None

    # 只有在上传新文件时，才执行转换
    if uploaded_file:
        # 创建临时目录存放文件
        temp_dir = tempfile.mkdtemp()

        # 将上传的 PDF 文件保存到临时目录
        pdf_path = os.path.join(temp_dir, uploaded_file.name)
        with open(pdf_path, "wb") as f:
            f.write(uploaded_file.getbuffer())

        # 显示上传的文件名和文件大小
        st.write(f"已上传文件: {uploaded_file.name}")
        st.write(f"文件大小: {uploaded_file.size} 字节")

    # 开始进行 PDF 转换为 PPT 按钮
    if uploaded_file:
        if st.button("开始转换为 PPT"):
            with st.spinner('正在转换文件，请稍候...'):
                try:
                    # 将 PDF 转换为图像
                    images = pdf_to_images(pdf_path)

                    # 生成 PPT 文件路径
                    ppt_filename = os.path.splitext(uploaded_file.name)[0] + ".pptx"
                    ppt_path = os.path.join(temp_dir, ppt_filename)

                    # 将图像插入到 PPT 文件中
                    images_to_ppt(images, ppt_path)

                    # 将处理结果保存在 session_state 中
                    st.session_state[uploaded_file.name] = {
                        "ppt_path1": ppt_path,
                        "ppt_filename1": ppt_filename,
                        "temp_dir1": temp_dir  # 保存临时目录路径
                    }

                    st.success("转换完成！")

                except Exception as e:
                    st.error(f"转换失败: {e}")

    # 显示下载按钮（如果转换完成）
    try:
        if uploaded_file and uploaded_file.name in st.session_state and st.session_state[uploaded_file.name] is not None:
            result = st.session_state[uploaded_file.name]
            ppt_path = result["ppt_path1"]
            ppt_filename = result["ppt_filename1"]

            # 提供 PPT 文件下载按钮
            with open(ppt_path, "rb") as ppt_file:
                st.download_button("下载转换后的 PPT 文件", ppt_file, file_name=ppt_filename)

            # 清理临时文件夹和文件
            if st.button('清理临时文件'):
                delete_temp_files(result["temp_dir1"])
                st.write("临时文件夹已删除。")
    except:
        pass


main()

